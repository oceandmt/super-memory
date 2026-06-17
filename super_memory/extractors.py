"""Document extractors for super_memory train pipeline."""
from __future__ import annotations

import importlib.util
from pathlib import Path
from typing import Any


def extract_pdf(path: Path) -> str | None:
    """Extract text from PDF using PyPDF2."""
    try:
        import PyPDF2
    except ImportError:
        return None
    
    try:
        with path.open("rb") as fh:
            reader = PyPDF2.PdfReader(fh)
            pages = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    pages.append(text)
            return "\n\n".join(pages)
    except Exception:
        return None


def extract_docx(path: Path) -> str | None:
    """Extract text from DOCX using python-docx."""
    try:
        import docx
    except ImportError:
        return None
    
    try:
        doc = docx.Document(path)
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        return "\n\n".join(paragraphs)
    except Exception:
        return None


def extract_pptx(path: Path) -> str | None:
    """Extract text from PPTX using python-pptx."""
    try:
        import pptx
    except ImportError:
        return None
    
    try:
        presentation = pptx.Presentation(path)
        slides = []
        for slide in presentation.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            if slide_text:
                slides.append("\n".join(slide_text))
        return "\n\n---\n\n".join(slides)
    except Exception:
        return None


def extract_html(path: Path) -> str | None:
    """Extract text from HTML using BeautifulSoup4."""
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        return None
    
    try:
        html = path.read_text(encoding="utf-8", errors="ignore")
        soup = BeautifulSoup(html, "html.parser")
        
        # Remove script/style tags
        for tag in soup(["script", "style", "nav", "header", "footer"]):
            tag.decompose()
        
        # Get text
        text = soup.get_text(separator="\n\n")
        
        # Clean whitespace
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        return "\n".join(lines)
    except Exception:
        return None


def extract_xlsx(path: Path) -> str | None:
    """Extract text from XLSX using openpyxl."""
    try:
        import openpyxl
    except ImportError:
        return None
    
    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        sheets = []
        for sheet in wb.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    rows.append(row_text)
            if rows:
                sheets.append(f"## {sheet.title}\n\n" + "\n".join(rows))
        return "\n\n".join(sheets)
    except Exception:
        return None


def extract_csv(path: Path) -> str | None:
    """Extract text from CSV."""
    try:
        import csv
    except ImportError:
        return None
    
    try:
        with path.open("r", encoding="utf-8", errors="ignore") as fh:
            reader = csv.reader(fh)
            rows = ["\t".join(row) for row in reader]
            return "\n".join(rows)
    except Exception:
        return None


EXTRACTORS: dict[str, Any] = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".html": extract_html,
    ".htm": extract_html,
    ".xlsx": extract_xlsx,
    ".csv": extract_csv,
}


def extract_text(path: Path) -> str | None:
    """
    Extract text from a file using appropriate extractor.
    
    Returns None if extractor not available or extraction failed.
    """
    ext = path.suffix.lower()
    
    # Plain text
    if ext in {".md", ".markdown", ".txt", ".rst"}:
        try:
            return path.read_text(encoding="utf-8", errors="ignore")
        except Exception:
            return None
    
    # Rich formats
    extractor = EXTRACTORS.get(ext)
    if extractor:
        return extractor(path)
    
    return None


def available_extractors() -> dict[str, bool]:
    """Check which extractors are available."""
    return {
        "pdf": importlib.util.find_spec("PyPDF2") is not None,
        "docx": importlib.util.find_spec("docx") is not None,
        "pptx": importlib.util.find_spec("pptx") is not None,
        "html": importlib.util.find_spec("bs4") is not None,
        "xlsx": importlib.util.find_spec("openpyxl") is not None,
        "csv": True,  # stdlib
        "text": True,  # builtin
    }
